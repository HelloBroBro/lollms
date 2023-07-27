from lollms.personality import APScript
from lollms.helpers import ASCIIColors, trace_exception
from lollms.paths import LollmsPaths
from sklearn.feature_extraction.text import TfidfVectorizer
import numpy as np
from pathlib import Path
import json

class TFIDFLoader:
    @staticmethod
    def create_vectorizer_from_dict(tfidf_info):
        vectorizer = TfidfVectorizer(**tfidf_info['params'])
        vectorizer.vocabulary_ = tfidf_info['vocabulary']
        vectorizer.idf_ = [tfidf_info['idf_values'][feature] for feature in vectorizer.get_feature_names()]
        return vectorizer

    @staticmethod
    def create_dict_from_vectorizer(vectorizer):
        tfidf_info = {
            "vocabulary": vectorizer.vocabulary_,
            "idf_values": dict(zip(vectorizer.get_feature_names(), vectorizer.idf_)),
            "params": vectorizer.get_params()
        }
        return tfidf_info
class TextVectorizer:
    def __init__(
                    self, 
                    vectorization_method, # supported "model_embedding" or "ftidf_vectorizer"
                    model=None, #needed in case of using model_embedding
                    database_path=None,
                    save_db=False,
                    visualize_data_at_startup=False,
                    visualize_data_at_add_file=False,
                    visualize_data_at_generate=False,
                    data_visualization_method="PCA",
                    database_dict=None
                    ):
        
        self.vectorization_method = vectorization_method
        self.save_db = save_db
        self.model = model
        self.database_file = database_path

        self.visualize_data_at_startup=visualize_data_at_startup
        self.visualize_data_at_add_file=visualize_data_at_add_file
        self.visualize_data_at_generate=visualize_data_at_generate
        
        self.data_visualization_method = data_visualization_method
        
        if database_dict is not None:
            self.chunks = []
            self.embeddings = database_dict["embeddings"]
            self.texts =  database_dict["text"]
            self.infos =   database_dict["infos"]
            self.ready = True
            self.vectorizer = database_dict["vectorizer"]
        else:
            self.chunks = []
            self.embeddings = {}
            self.texts = {}
            self.ready = False
            self.vectorizer = None
        
            if vectorization_method=="model_embedding":
                try:
                    if not self.model or self.model.embed("hi")==None: # test
                        self.vectorization_method="ftidf_vectorizer"
                        self.infos={
                            "vectorization_method":"ftidf_vectorizer"
                        }
                    else:
                        self.infos={
                            "vectorization_method":"model_embedding"
                        }
                except Exception as ex:
                    ASCIIColors.error("Couldn't embed the text, so trying to use tfidf instead.")
                    trace_exception(ex)
                    self.infos={
                        "vectorization_method":"ftidf_vectorizer"
                    }
        # Load previous state from the JSON file
        if self.save_db:
            if Path(self.database_file).exists():
                ASCIIColors.success(f"Database file found : {self.database_file}")
                self.load_from_json()
                if self.visualize_data_at_startup:
                    self.show_document()
                self.ready = True
            else:
                ASCIIColors.info(f"No database file found : {self.database_file}")

                
    def show_document(self, query_text=None, save_fig_path =None, show_interactive_form=False):
        import textwrap
        import seaborn as sns
        import matplotlib.pyplot as plt
        import mplcursors
        from tkinter import Tk, Text, Scrollbar, Frame, Label, TOP, BOTH, RIGHT, LEFT, Y, N, END

        
        from sklearn.manifold import TSNE
        from sklearn.decomposition import PCA

        if self.data_visualization_method=="PCA":
            use_pca =  True
        else:
            use_pca =  False
        
        if use_pca:
            print("Showing pca representation :")
        else:
            print("Showing t-sne representation :")
        texts = list(self.texts.values())
        embeddings = self.embeddings
        emb = list(embeddings.values())
        ref = list(embeddings.keys())
        if len(emb)>=2:
            # Normalize embeddings
            emb = np.vstack(emb)
            norms = np.linalg.norm(emb, axis=1)
            normalized_embeddings = emb / norms[:, np.newaxis]

            # Embed the query text
            if query_text is not None:
                query_embedding = self.embed_query(query_text)
                query_embedding = query_embedding.detach().squeeze().numpy()
                query_normalized_embedding = query_embedding / np.linalg.norm(query_embedding)

                # Combine the query embedding with the document embeddings
                combined_embeddings = np.vstack((normalized_embeddings, query_normalized_embedding))
                ref.append("Quey_chunk_0")
            else:
                # Combine the query embedding with the document embeddings
                combined_embeddings = normalized_embeddings

            if use_pca:
                # Use PCA for dimensionality reduction
                pca = PCA(n_components=2)
                try:
                    embeddings_2d = pca.fit_transform(combined_embeddings)
                except Exception as ex:
                    
                    embeddings_2d = []
            else:
                # Use t-SNE for dimensionality reduction
                # Adjust the perplexity value
                perplexity = min(30, combined_embeddings.shape[0] - 1)
                tsne = TSNE(n_components=2, perplexity=perplexity)
                embeddings_2d = tsne.fit_transform(combined_embeddings)

            # Create a dictionary to map document paths to colors
            document_path_colors = {}
            for i, path in enumerate(ref):
                document_path = "_".join(path.split("_")[:-1])  # Extract the document path (excluding chunk and chunk number)
                if document_path not in document_path_colors:
                    # Assign a new color to the document path if it's not in the dictionary
                    document_path_colors[document_path] = sns.color_palette("hls", len(document_path_colors) + 1)[-1]

            # Generate a list of colors for each data point based on the document path
            point_colors = [document_path_colors["_".join(path.split("_")[:-1])] for path in ref]

            # Create a scatter plot using Seaborn
            sns.scatterplot(x=embeddings_2d[:, 0], y=embeddings_2d[:, 1], hue=point_colors)  # Plot document embeddings
            # Add labels to the scatter plot
            for i, (x, y) in enumerate(embeddings_2d[:-1]):
                plt.text(x, y, str(i), fontsize=8)

            plt.xlabel('Dimension 1')
            plt.ylabel('Dimension 2')
            if use_pca:      
                plt.title('Embeddings Scatter Plot based on PCA')
            else:
                plt.title('Embeddings Scatter Plot based on t-SNE')
            # Enable mplcursors to show tooltips on hover
            cursor = mplcursors.cursor(hover=True)

            # Define the hover event handler
            @cursor.connect("add")
            def on_hover(sel):
                index = sel.target.index
                if index > 0:
                    text = texts[index]
                    wrapped_text = textwrap.fill(text, width=50)  # Wrap the text into multiple lines
                    sel.annotation.set_text(f"Index: {index}\nText:\n{wrapped_text}")
                else:
                    sel.annotation.set_text("Query")

            # Define the click event handler using matplotlib event handling mechanism
            def on_click(event):
                if event.xdata is not None and event.ydata is not None:
                    x, y = event.xdata, event.ydata
                    distances = ((embeddings_2d[:, 0] - x) ** 2 + (embeddings_2d[:, 1] - y) ** 2)
                    index = distances.argmin()
                    text = texts[index] if index < len(texts) else query_text

                    # Open a new Tkinter window with the content of the text
                    root = Tk()
                    root.title(f"Text for Index {index}")
                    frame = Frame(root)
                    frame.pack(fill=BOTH, expand=True)

                    label = Label(frame, text="Text:")
                    label.pack(side=TOP, padx=5, pady=5)

                    text_box = Text(frame)
                    text_box.pack(side=TOP, padx=5, pady=5, fill=BOTH, expand=True)
                    text_box.insert(END, text)

                    scrollbar = Scrollbar(frame)
                    scrollbar.pack(side=RIGHT, fill=Y)
                    scrollbar.config(command=text_box.yview)
                    text_box.config(yscrollcommand=scrollbar.set)

                    text_box.config(state="disabled")

                    root.mainloop()

            # Connect the click event handler to the figure
            plt.gcf().canvas.mpl_connect("button_press_event", on_click)
            if save_fig_path:
                plt.savefig(save_fig_path)
            if show_interactive_form:
                plt.show()
        
    def add_document(self, document_id, text, chunk_size, overlap_size, force_vectorize=False):
        if document_id in self.embeddings and not force_vectorize:
            print(f"Document {document_id} already exists. Skipping vectorization.")
            return

        # Split tokens into sentences
        sentences = text.split('. ')
        def remove_empty_sentences(sentences):
            return [self.model.tokenize(sentence) for sentence in sentences if sentence.strip() != '']
        sentences = remove_empty_sentences(sentences)
        # Generate chunks with overlap and sentence boundaries
        chunks = []
        current_chunk = []
        for i in range(len(sentences)):
            sentence_tokens = sentences[i]

            # ASCIIColors.yellow(len(sentence_tokens))
            if len(current_chunk) + len(sentence_tokens) <= chunk_size:
                current_chunk.extend(sentence_tokens)
            else:
                if current_chunk:
                    chunks.append(current_chunk)
                    
                current_chunk=[]
                for j in reversed(range(overlap_size)):
                    current_chunk.extend(sentences[i-j-1])
                current_chunk.extend(sentence_tokens)
            

        if current_chunk:
            for i, chunk_text in enumerate(chunks):
                chunk_id = f"{document_id}_chunk_{i + 1}"
                chunk_dict = {
                    "chunk_id": chunk_id,
                    "chunk_text": chunk_text
                }
                self.chunks.append(chunk_dict)
        
    def index(self):
        if self.vectorization_method=="ftidf_vectorizer":
            self.vectorizer = TfidfVectorizer()
            #if self.debug:
            #    ASCIIColors.yellow(','.join([len(chunk) for chunk in chunks]))
            data=[]
            for chunk in self.chunks:
                try:
                    data.append(self.model.detokenize(chunk["chunk_text"]).replace("<s>","").replace("</s>","") ) 
                except Exception as ex:
                    print("oups")
            self.vectorizer.fit(data)

        self.embeddings = {}
        # Generate embeddings for each chunk
        for i, chunk in enumerate(self.chunks):
            # Store chunk ID, embedding, and original text
            chunk_id = chunk["chunk_id"]
            chunk_text = chunk["chunk_text"]
            try:
                self.texts[chunk_id] = self.model.detokenize(chunk_text)
                if self.vectorization_method=="ftidf_vectorizer":
                    self.embeddings[chunk_id] = self.vectorizer.transform([self.texts[chunk_id]]).toarray()
                else:
                    self.embeddings[chunk_id] = self.model.embed(self.texts[chunk_id])
            except Exception as ex:
                print("oups")

        if self.save_db:
            self.save_to_json()
            
        self.ready = True
        if self.visualize_data_at_add_file:
            self.show_document()


    def embed_query(self, query_text):
        # Generate query embedding
        if self.vectorization_method=="ftidf_vectorizer":
            query_embedding = self.vectorizer.transform([query_text]).toarray()
        else:
            query_embedding = self.model.embed(query_text)
            if query_embedding is None:
                ASCIIColors.warning("The model doesn't implement embedding extraction")
                self.vectorization_method="ftidf_vectorizer"
                query_embedding = self.vectorizer.transform([query_text]).toarray()

        return query_embedding

    def recover_text(self, query_embedding, top_k=1):
        from sklearn.metrics.pairwise import cosine_similarity
        similarities = {}
        for chunk_id, chunk_embedding in self.embeddings.items():
            similarity = cosine_similarity(query_embedding, chunk_embedding)
            similarities[chunk_id] = similarity

        # Sort the similarities and retrieve the top-k most similar embeddings
        sorted_similarities = sorted(similarities.items(), key=lambda x: x[1], reverse=True)[:top_k]

        # Retrieve the original text associated with the most similar embeddings
        texts = [self.texts[chunk_id] for chunk_id, _ in sorted_similarities]

        if self.visualize_data_at_generate:
            self.show_document()

        return texts, sorted_similarities

    def toJson(self):
        state = {
            "embeddings": {str(k): v.tolist()  if type(v)!=list else v for k, v in self.embeddings.items() },
            "texts": self.texts,
            "infos": self.infos,
            "vectorizer": TFIDFLoader.create_vectorizer_from_dict(self.vectorizer) if self.vectorization_method=="ftidf_vectorizer" else None
        }
        return state
    
    def setVectorizer(self, vectorizer_dict:dict):
        self.vectorizer=TFIDFLoader.create_vectorizer_from_dict(vectorizer_dict)

    def save_to_json(self):
        state = {
            "embeddings": {str(k): v.tolist()  if type(v)!=list else v for k, v in self.embeddings.items() },
            "texts": self.texts,
            "infos": self.infos
        }
        with open(self.database_file, "w") as f:
            json.dump(state, f)

    def load_from_json(self):

        ASCIIColors.info("Loading vectorized documents")
        with open(self.database_file, "r") as f:
            state = json.load(f)
            self.embeddings = {k: v for k, v in state["embeddings"].items()}
            self.texts = state["texts"]
            self.infos= state["infos"]
            self.ready = True
        if self.vectorization_method=="ftidf_vectorizer":
            from sklearn.feature_extraction.text import TfidfVectorizer
            data = list(self.texts.values())
            if len(data)>0:
                self.vectorizer = TfidfVectorizer()
                self.vectorizer.fit(data)
                self.embeddings={}
                for k,v in self.texts.items():
                    self.embeddings[k]= self.vectorizer.transform([v]).toarray()
                    
                    
    def clear_database(self):
        self.vectorizer=None
        self.embeddings = {}
        self.texts={}
        if self.save_db:
            self.save_to_json()
            
      
class GenericDataLoader:
    @staticmethod
    def install_package(package_name):
        import subprocess
        import sys
        subprocess.check_call([sys.executable, "-m", "pip", "install", package_name])

    @staticmethod        
    def read_pdf_file(file_path):
        try:
            import PyPDF2
        except ImportError:
            GenericDataLoader.install_package("PyPDF2")
            import PyPDF2
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
        return text

    @staticmethod
    def read_docx_file(file_path):
        try:
            from docx import Document
        except ImportError:
            GenericDataLoader.install_package("python-docx")
            from docx import Document
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text

    @staticmethod
    def read_json_file(file_path):
        import json
        with open(file_path, 'r') as file:
            data = json.load(file)
        return data
    
    @staticmethod
    def read_csv_file(file_path):
        try:
            import csv
        except ImportError:
            GenericDataLoader.install_package("csv")
            import csv
        with open(file_path, 'r') as file:
            csv_reader = csv.reader(file)
            lines = [row for row in csv_reader]
        return lines    

    @staticmethod
    def read_html_file(file_path):
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            GenericDataLoader.install_package("beautifulsoup4")
            from bs4 import BeautifulSoup
        with open(file_path, 'r') as file:
            soup = BeautifulSoup(file, 'html.parser')
            text = soup.get_text()
        return text
    
    @staticmethod
    def read_pptx_file(file_path):
        try:
            from pptx import Presentation
        except ImportError:
            GenericDataLoader.install_package("python-pptx")
            from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
        return text
    
    @staticmethod
    def read_text_file(file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        return content